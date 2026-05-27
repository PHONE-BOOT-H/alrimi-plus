"""출품작 PPTX 생성. python -m scripts.build_pptx

build_deck.py가 만든 _deck.html을 2x 고해상도로 렌더 → 15슬라이드(16:9) PPTX.
PDF와 픽셀 단위로 동일한 품질. 먼저 build_deck.py를 실행해 _deck.html이 있어야 함.
"""
from __future__ import annotations

import os
import tempfile

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

SUB = r"C:\Users\hanta\Desktop\한태영\대외활동\공모전\교육공공데이터활용대회\14_제출본"
HTML = os.path.join(SUB, "_deck.html")
OUT = os.path.join(SUB, "알리미플러스_출품작.pptx")

tmp = tempfile.mkdtemp(prefix="alrimi_pptx_")
imgs: list[str] = []
with sync_playwright() as p:
    br = p.chromium.launch()
    pg = br.new_page(viewport={"width": 1280, "height": 720}, device_scale_factor=2)
    pg.goto("file:///" + HTML.replace("\\", "/"), wait_until="networkidle")
    pg.evaluate("document.fonts.ready")
    pg.wait_for_timeout(1800)
    pages = pg.query_selector_all(".page")
    for i, el in enumerate(pages):
        path = os.path.join(tmp, f"s{i+1:02d}.png")
        el.screenshot(path=path)
        imgs.append(path)
    br.close()

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
for img in imgs:
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)
prs.save(OUT)
print("PPTX:", len(imgs), "slides,", round(os.path.getsize(OUT) / 1024 / 1024, 2), "MB")
